import os
import json
import requests
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from groq import Groq
from dotenv import load_dotenv

# Load environment variables from .env
load_dotenv()

# ---------------------- Theme Definitions ---------------------- #
THEMES = {
    "Modern Blue": {
        "bg_color": RGBColor(0, 102, 204),
        "title_font": ("Calibri", 44, RGBColor(255, 255, 255)),
        "content_font": ("Calibri", 24, RGBColor(255, 255, 255)),
        "accent_color": RGBColor(255, 215, 0)  # gold bullets
    },
    "Classic Dark": {
        "bg_color": RGBColor(34, 34, 34),
        "title_font": ("Times New Roman", 44, RGBColor(255, 255, 255)),
        "content_font": ("Times New Roman", 24, RGBColor(200, 200, 200)),
        "accent_color": RGBColor(0, 255, 0)  # green bullets
    },
    "Minimal White": {
        "bg_color": RGBColor(255, 255, 255),
        "title_font": ("Arial", 44, RGBColor(0, 0, 0)),
        "content_font": ("Arial", 24, RGBColor(51, 51, 51)),
        "accent_color": RGBColor(0, 102, 204)
    }
}

# ---------------------- PPT Generator ---------------------- #
class PPTGenerator:
    def __init__(self, groq_api_key, theme="Modern Blue"):
        if not groq_api_key:
            raise ValueError("Groq API key is required.")
        self.client = Groq(api_key=groq_api_key)
        self.presentation = Presentation()
        self.theme = theme

    def generate_content_outline(self, topic, num_slides=5):
        prompt = f"""
        Generate a PowerPoint outline for topic: "{topic}" with {num_slides} slides.
        Return JSON array like:
        [
            {{
                "title": "Slide Title",
                "content": ["Bullet point 1", "Bullet point 2"],
                "slide_type": "title|content|image|conclusion",
                "image_prompt": "Optional image description"
            }}
        ]
        """
        try:
            response = self.client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=1500  # Increased for more content
            )
            text = response.choices[0].message.content.strip()
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0].strip()
            elif "```" in text:
                text = text.split("```")[1].strip()
            slides = json.loads(text)
            return slides
        except Exception as e:
            print(f"Error generating outline: {e}")
            return self._fallback_outline(topic, num_slides)

    def _fallback_outline(self, topic, num_slides):
        return [
            {"title": f"Introduction to {topic}", "content": ["Overview", "Objectives"], "slide_type": "title", "image_prompt": ""},
            {"title": "Key Concepts", "content": ["Main principles", "Core ideas"], "slide_type": "content", "image_prompt": ""},
            {"title": "Applications", "content": ["Examples", "Case studies"], "slide_type": "content", "image_prompt": ""},
            {"title": "Challenges", "content": ["Limitations", "Ethics"], "slide_type": "content", "image_prompt": ""},
            {"title": "Conclusion", "content": ["Summary", "Next steps"], "slide_type": "conclusion", "image_prompt": ""}
        ]

    def download_image(self, query, save_path="temp_image.jpg", unsplash_api_key=None):
        if not unsplash_api_key:
            print("⚠️ UNSPLASH API key not provided. Using placeholder image.")
            return self._placeholder_image(save_path)
        try:
            url = "https://api.unsplash.com/photos/random"
            params = {"query": query, "client_id": unsplash_api_key}
            response = requests.get(url, params=params)
            response.raise_for_status()
            data = response.json()
            image_url = data["urls"]["regular"]
            img_response = requests.get(image_url)
            with open(save_path, "wb") as f:
                f.write(img_response.content)
            return save_path
        except Exception as e:
            print(f"Error downloading image: {e}")
            return self._placeholder_image(save_path)

    def _placeholder_image(self, save_path):
        img = Image.new("RGB", (800, 600), color="#4A90E2")
        img.save(save_path)
        return save_path

    # ------------------ Apply Theme ------------------ #
    def apply_theme_to_slide(self, slide, is_content=True):
        theme = THEMES.get(self.theme, THEMES["Modern Blue"])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = theme["bg_color"]

        # Title formatting
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text_frame.paragraphs[0].font.name = theme["title_font"][0]
            title_shape.text_frame.paragraphs[0].font.size = Pt(theme["title_font"][1])
            title_shape.text_frame.paragraphs[0].font.color.rgb = theme["title_font"][2]

        # Content formatting
        if is_content:
            for shape in slide.shapes:
                if shape.has_text_frame and shape != title_shape:
                    for p in shape.text_frame.paragraphs:
                        p.font.name = theme["content_font"][0]
                        p.font.size = Pt(theme["content_font"][1])
                        p.font.color.rgb = theme["accent_color"]  # Accent for bullets

    # ------------------ Slide Creation ------------------ #
    def create_title_slide(self, title, subtitle=""):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
        self.apply_theme_to_slide(slide, is_content=False)
        return slide

    def create_content_slide(self, title, content, include_image=False, image_prompt="", unsplash_api_key=None):
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        content_text = "\n".join([f"• {c}" for c in content]) if isinstance(content, list) else str(content)
        slide.placeholders[1].text = content_text
        self.apply_theme_to_slide(slide)
        if include_image and image_prompt:
            path = self.download_image(image_prompt, unsplash_api_key=unsplash_api_key)
            if path:
                slide.shapes.add_picture(path, Inches(5), Inches(2), height=Inches(4))
                os.remove(path)
        return slide

    def create_image_slide(self, title, content, image_prompt, unsplash_api_key=None):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_box.text_frame.text = title
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
        content_box.text_frame.text = "\n".join([f"• {c}" for c in content]) if isinstance(content, list) else str(content)
        self.apply_theme_to_slide(slide)
        path = self.download_image(image_prompt, unsplash_api_key=unsplash_api_key)
        if path:
            slide.shapes.add_picture(path, Inches(5), Inches(2), height=Inches(5))
            os.remove(path)
        return slide

    # ------------------ Generate Presentation ------------------ #
    def generate_presentation(self, topic, num_slides=5, output_path="generated_presentation.pptx", unsplash_api_key=None):
        outline = self.generate_content_outline(topic, num_slides)
        for i, slide_data in enumerate(outline, start=1):
            title = slide_data["title"]
            content = slide_data["content"]
            slide_type = slide_data.get("slide_type", "content")
            image_prompt = slide_data.get("image_prompt", "")
            if i == 1 or slide_type == "title":
                self.create_title_slide(title, "Generated by Groq AI")
            elif slide_type == "image":
                self.create_image_slide(title, content, image_prompt, unsplash_api_key)
            else:
                include_image = (i % 2 == 0)
                self.create_content_slide(title, content, include_image, image_prompt, unsplash_api_key)
        self.presentation.save(output_path)
        return output_path

# ---------------------- Streamlit App ---------------------- #
def main():
    st.title("📊 AI PowerPoint Generator with Themes")
    st.write("Generate AI-powered PowerPoint presentations using Groq + Unsplash")

    topic = st.text_input("Enter Topic", "Cricket and Its History")
    num_slides = st.slider("Number of Slides", 3, 15, 5)
    theme = st.selectbox("Select Theme", list(THEMES.keys()))
    
    groq_api_key = os.getenv("GROQ_API_KEY")
    unsplash_api_key = os.getenv("UNSPLASH_API_KEY")

    if st.button("Generate Presentation"):
        if not groq_api_key:
            st.error("Please set your GROQ_API_KEY in .env file!")
            return
        with st.spinner("⏳ Generating presentation..."):
            ppt = PPTGenerator(groq_api_key=groq_api_key, theme=theme)
            output_file = ppt.generate_presentation(topic, num_slides, unsplash_api_key=unsplash_api_key)
            with open(output_file, "rb") as f:
                st.success("✅ Presentation Generated!")
                st.download_button(
                    label="📥 Download PPT",
                    data=f,
                    file_name=output_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

if __name__ == "__main__":
    main()
